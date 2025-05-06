import base64
import json
from io import BytesIO
from flask import Flask, request, jsonify
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
import openpyxl

app = Flask(__name__)

def extract_pdf_text(file_content):
    reader = PdfReader(BytesIO(file_content))
    text = ""
    for page in reader.pages:
        text += page.extract_text()
    return text

def extract_pptx_text(file_content):
    presentation = Presentation(BytesIO(file_content))
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

def extract_docx_text(file_content):
    document = Document(BytesIO(file_content))
    text = ""
    for para in document.paragraphs:
        text += para.text
    return text

def extract_xlsx_text(file_content):
    workbook = openpyxl.load_workbook(BytesIO(file_content))
    text = ""
    for sheet in workbook.sheetnames:
        worksheet = workbook[sheet]
        for row in worksheet.iter_rows(values_only=True):
            for cell in row:
                if cell:
                    text += str(cell)
    return text

@app.route('/extract', methods=['POST'])
def extract_text():
    data = request.get_json()
    file_name = data['fileName']
    file_content = base64.b64decode(data['fileContent'])

    if file_name.lower().endswith('.pdf'):
        text = extract_pdf_text(file_content)
    elif file_name.lower().endswith('.pptx'):
        text = extract_pptx_text(file_content)
    elif file_name.lower().endswith('.docx'):
        text = extract_docx_text(file_content)
    elif file_name.lower().endswith('.xlsx'):
        text = extract_xlsx_text(file_content)
    else:
        return jsonify({"error": "Unsupported file type"}), 400

    return jsonify({"text": text})

if __name__ == '__main__':
    app.run(debug=True)
